import os
import simplify
import tarfile
import shutil
import re

from get_lyrics import requestsScrape, divide_by_text
import simplify

import matplotlib.colors
from PIL import Image, ImageDraw, ImageFont

from pptx import Presentation
from pptx.util import Pt, Inches, Cm
from pptx.util import Length
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN


# def find(s, ch):
#     return (i for i, ltr in enumerate(s) if ltr == ch)


def split_sentence(sentence, times):
    """Splits the sentence in ``times`` parts"""
    idxs = [int((i+1) * (len(sentence) / times)) for i in range(times-1)]
    spc_idxs = [0]
    for i, char in enumerate(sentence):
        if char == " ": 
            spc_idxs.append(i)
    spc_idxs.append(None)

    splited = [] 

    if len(spc_idxs) - 2 < times:
        print(sentence, times)
        raise Exception("\"times\" must be less or equal the number of blank spaces!")
    elif len(spc_idxs) - 2 == times:
        for i, idx in enumerate(spc_idxs):
            if idx == None:
                break
            next = spc_idxs[i+1]
            splited.append(sentence[idx:next])
    else:
        difference = 100
        closest_idx = 0
        idxs_mapping = [0]

        spc_idxs[-1] = len(spc_idxs)

        for i, idx in enumerate(idxs):
            for j, spc_idx in enumerate(spc_idxs):
                if abs(spc_idx - idx) < difference:
                    difference = abs(spc_idx - idx)
                    closest_idx = spc_idx
            idxs_mapping.append(closest_idx)
            spc_idxs.remove(closest_idx)
            difference = 100

        idxs_mapping = sorted(idxs_mapping)
        idxs_mapping.append(None)

        for i, idx in enumerate(idxs_mapping):
            if idx == None:
                break
            next = idxs_mapping[i+1]
            splited.append(sentence[idx:next])
            
    return [i.strip() for i in splited]


class Px(Length):
    """
    Convenience value class for specifying a length in pixels
    """

    def __new__(cls, px):
        return Length.__new__(cls, Cm(px / 118.121554))


class Colors(object):
    """
    Its attributes are color names and their values are rgb tuples
    """
    def __init__(self):
        for name, hex in matplotlib.colors.cnames.items():
            color = tuple([int(i * 255) for i in matplotlib.colors.to_rgb(hex)])
            setattr(self, name, color)


class Positions(object):
    """
    Maps all possible positions
    """
    def __init__(self, width, height, w, h, shadow=False):
        self.width = width
        self.height = height
        self.w = w
        self.h = h

        self.shadow = shadow

        self.border = 80
        self.extra = 5

        self.update()
    
    def update(self):
        top = self.border + self.extra * self.shadow
        bottom = self.height - self.border - self.h + self.extra * self.shadow
        h_center = int((self.height - self.h) / 2) + self.extra * self.shadow
        w_center = int((self.width - self.w) / 2) - self.extra * self.shadow
        right = self.width - self.w - self.border * (self.width / self.height) - self.extra * self.shadow
        left = self.border * (self.width / self.height) - self.extra * self.shadow

        self.middle = (w_center, h_center)
        self.center_left = (left, h_center)
        self.center_right = (right, h_center)

        self.bottom_left = (left, bottom)
        self.bottom_right = (right, bottom)
        self.bottom_center = (w_center, bottom)

        self.top_left = (left, top)
        self.top_right = (right, top)
        self.top_center = (w_center, top)

    
    @staticmethod
    def get_font_size(text, font_type):
        if "\n" in text:
            greater = 0
            h = 0
            for sentence in text.split("\n"):
                size = font_type.getsize(sentence)
                w = size[0]
                greater = w if w > greater else greater

                h += font_type.getsize("A")[1] + 2

            w = greater
        else:
            w, h = font_type.getsize(text)

        return w, h
    

    @staticmethod
    def check_alignment(font_name):
        translate_alignment = {
            'middle': 'center',
            'center_left': 'left',
            'center_right': 'right',
            'bottom_left': 'left',
            'bottom_right': 'right',
            'bottom_center': 'center',
            'top_left': 'left',
            'top_right': 'right',
            'top_center': 'center',
        }

        if font_name in translate_alignment:
            return translate_alignment[font_name]
        else:
            raise Exception("Invalid position!")


class FontType(object):
    def __init__(self, family_font, size, weigth=False):
        self._family_font = family_font
        self._size = size
        self._weigth = weigth

    def get_path(self):
        return simplify.get_font_path(self._family_font, "bold" if self._weigth else None)
    
    def get_size(self):
        return self._size
    
    def is_bold(self):
        return self._weigth

    def set_size(self, size):
        self._size = size
    
    def set_family(self, family):
        self._family_font = family
    
    def set_weight(self, weight):
        self._weigth = weight


class Slide(object):
    def __init__(self):
        self.width = 1920
        self.height = 1080

        self.font_type = FontType("Ubuntu", 40, "bold")

        self.color = getattr(Colors(), "white")

        self.stacks = 2 # Number of sentences per slide

        self.directory = "../slides/Slides"

        self.border = 20

        self.extra = 5 # Shadow distance

        self.position = "middle"
    

    def adapt_size(self, lyrics, limit):
        """
        Returns a new list with the verses that are larger then the limit splited
        """
        new = []
        for i, estrofe in enumerate(lyrics):
            new.append([])
            for j, verso in enumerate(estrofe):
                font_type = ImageFont.truetype(self.font_type.get_path(), self.font_type.get_size())
                text_width = Positions.get_font_size(verso, font_type)[0]

                if text_width > limit:
                    split_times = text_width // limit

                    for subverso in split_sentence(verso, split_times):
                        new[i].append(subverso)
                else:
                    new[i].append(verso)
        
        return new
    
    def create_imageshow(self, background_path, lyrics, shadow=False):
        """
        Creates an image show
        """
        image = simplify.assign_image(background_path)

        os.mkdir(self.directory)

        font_type = ImageFont.truetype(self.font_type.get_path(), self.font_type.get_size())

        normal_positions = None
        shadow_positions = None

        count = 1
        for estrofe in lyrics:
            for i in range(0, len(estrofe), self.stacks):
                if i < len(estrofe) - (self.stacks - 1):
                    text = "\n".join([estrofe[j] for j in range(i, i+self.stacks)])
                else:
                    last = len(estrofe)
                    text = "\n".join([estrofe[j] for j in range(i, last)])
                
                w, h = Positions.get_font_size(text, font_type)

                if not normal_positions or not shadow_positions:
                    normal_positions = Positions(self.width, self.height, w, h)
                    shadow_positions = Positions(self.width, self.height, w, h, shadow=True)
                else:
                    normal_positions.w, normal_positions.h = w, h
                    shadow_positions.w, shadow_positions.h = w, h

                normal_positions.extra = self.extra
                normal_positions.border = self.border

                shadow_positions.extra = self.extra
                shadow_positions.border = self.border

                normal_positions.update()
                shadow_positions.update()

                normal_position = getattr(normal_positions, self.position)
                shadow_position = getattr(shadow_positions, self.position)

                image = simplify.assign_image(background_path)
                if shadow:
                    image = image.copy()
                    draw = ImageDraw.Draw(image)
                    draw.text(xy=shadow_position, text=text, fill=getattr(Colors(), "black"), \
                        font=font_type, align=Positions.check_alignment(self.position))

                image = image.copy()
                draw = ImageDraw.Draw(image)
                draw.text(xy=normal_position, text=text, fill=self.color, \
                    font=font_type, align=Positions.check_alignment(self.position))
                
                image.save(f"{self.directory}/{count}.png")

                count += 1
    
    def create_pwp(self, prs, background_path, text, shadow=False):
        """
        Creates a pptx presentation based on the music lyrics
        """

        font_type = ImageFont.truetype(self.font_type.get_path(), self.font_type.get_size())
        w, h = Positions.get_font_size(text, font_type)

        normal_positions = Positions(self.width, self.height, w, h)
        shadow_positions = Positions(self.width, self.height, w, h, shadow=True)

        normal_positions.extra = self.extra
        normal_positions.border = self.border
        
        shadow_positions.extra = self.extra
        shadow_positions.border = self.border

        normal_positions.update()
        shadow_positions.update()

        normal_position = getattr(normal_positions, self.position)
        shadow_position = getattr(shadow_positions, self.position)

        layout = prs.slide_layouts[6]
        
        slide = prs.slides.add_slide(layout)

        # slide.shapes._spTree.insert(0, picture._element)
        slide.shapes.add_picture(background_path, 0, 0, height = Px(self.height))

        alignment = Positions.check_alignment(self.position)
        
        if len(self.position.split("_")) > 1 and self.position.split("_")[0] == "top":
            border_plus = -70
            multiplier_top = 1
        elif len(self.position.split("_")) > 1 and self.position.split("_")[0] == "bottom":
            border_plus = 0
            multiplier_top = -1
        else:
            border_plus = 0
            multiplier_top = 0
        
        if len(self.position.split("_")) > 1 and self.position.split("_")[1] == "left":
            multiplier_left = -1
        elif len(self.position.split("_")) > 1 and self.position.split("_")[1] == "right":
            multiplier_left = 1
        else:
            multiplier_left = 0

        if shadow:
            txBox = slide.shapes.add_textbox(0, 0, Px(w), Px(h))

            tf = txBox.text_frame

            tf.word_wrap = True

            txBox.width = prs.slide_width
            txBox.height = prs.slide_height
            txBox.top = Px(shadow_position[1] + border_plus + self.border * multiplier_top)
            txBox.left = Px(0 - self.extra - self.border * multiplier_left * self.width / self.height)

            p = tf.paragraphs[0]
            p.text = text
            p.alignment = getattr(PP_ALIGN, alignment.upper())

            font = p.font

            font.name = self.font_type.get_path()
            font.size = Px(self.font_type.get_size())
            font.color.rgb = RGBColor(*getattr(Colors(), "black"))

            font.bold = True if self.font_type.is_bold() else False

        txBox = slide.shapes.add_textbox(0, 0, Px(w), Px(h))

        tf = txBox.text_frame

        tf.word_wrap = True

        txBox.width = prs.slide_width
        txBox.height = prs.slide_height
        txBox.top = Px(normal_position[1] + border_plus + self.border * multiplier_top)
        txBox.left = Px(0 - self.border * multiplier_left * self.width / self.height)

        p = tf.paragraphs[0]
        p.text = text
        p.alignment = getattr(PP_ALIGN, alignment.upper())

        font = p.font
        
        font.name = self.font_type.get_path()
        font.size = Px(self.font_type.get_size())
        font.color.rgb = RGBColor(*self.color)

        font.bold = True if self.font_type.is_bold() else False
    

    def create_slideshow(self, background_path, lyrics, shadow=False, font_format=[]):
        prs = Presentation()
        prs.slide_width = Px(self.width)
        prs.slide_height = Px(self.height)

        LIMIT = 450
        for i, estrofe in enumerate(self.adapt_size(lyrics, LIMIT)):
            for i in range(0, len(estrofe), self.stacks):
                if i < len(estrofe) - (self.stacks - 1):
                    text = "\n".join([estrofe[j].capitalize() for j in range(i, i+self.stacks)])
                else:
                    last = len(estrofe)
                    text = "\n".join([estrofe[j].capitalize() for j in range(i, last)])

                font_type = ImageFont.truetype(self.font_type.get_path(), self.font_type.get_size())

                self.create_pwp(prs, "/home/mateusap1/Pictures/Outros/Fundo Letra da Música.png",
                                text, shadow)

        prs.save(f"{self.directory}.pptx")


if __name__ == "__main__":
    font_size = 72
    slide_presentation = Slide()
    slide_presentation.font_type.set_size(font_size)
    slide_presentation.border = 72
    slide_presentation.extra = 10 * font_size / 72
    slide_presentation.stacks = 2
    lyrics = requestsScrape("Projeto Sola", "23")
    for position in ["middle", "center_left", "center_right", "top_left", \
                     "top_center", "top_right", "bottom_left", "bottom_center", "bottom_right"]:
        slide_presentation.position = position
        slide_presentation.directory = f"../slides/{position}"

        slide_presentation.create_slideshow("/home/mateusap1/Pictures/Outros/Fundo Letra da Música.png", \
                                            lyrics, True)
        slide_presentation.create_imageshow("/home/mateusap1/Pictures/Outros/Fundo Letra da Música.png", \
            lyrics, shadow=True)