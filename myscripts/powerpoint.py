from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN

from get_lyrics import requestsScrape, divide_by_text
from simplify import get_font_path
from image import compress

from PIL import ImageFont

import os
import time


# Currently the GUI doesn't have a powerpoint option because it would be really specific, but you can adapt yor program according to your needs.
# Atualmente a GUI não tem uma opção para escolher o modo "powerpoint" porque seria muito específico, mas você pode adaptar essa função às suas necessidades.
def main():
    data = [
        "Algo novo - Arena Louvor", 
        "Tu És Bom (You Are Good) - Vineyard",
        "Ele Vem - Gabriel Guedes",
        "Coração Igual Teu - Diante do Trono",
        "(Oceanos) Onde meus pés podem Falhar - Ana Nóbrega",
        "Viverei pra Ti - Bereana Louvor"
    ]

    prs, layout = init_slides()
    for element in data:
        music, artist = element.split(" - ")
        # if music == "Meu Prazer":
        #     with open("letra.txt", "r") as letra:
        #         create_main(prs, layout, "Gézi Monteiro", "Meu Prazer", "Ubuntu", 32, text=letra.read())
        prs = create_main(prs, layout, artist, music, "Ubuntu", 32, shadow=False)
    
    prs.save("./letras/Slides.pptx")


def create_main(prs, layout, artist, music, font_name, size, text=None, color=RGBColor(255, 255, 255), stacks=1, shadow=True):
    if not text:
        letra = requestsScrape(artist, music)
    else:
        letra = divide_by_text(text)

    if not letra:
        print("There are no lyrics")
        return prs

    new = []

    LIMIT = 320
    for i, estrofe in enumerate(letra):
        new.append([])
        for j, verso in enumerate(estrofe):
            if ImageFont.truetype(get_font_path(font_name), size).getsize(verso)[0] > LIMIT:
                spc_idx = 0
                for n, char in enumerate(verso):
                    if char == " " and (abs(n - len(verso) // 2) < abs(spc_idx - len(verso) // 2)):
                        spc_idx = n
                
                new[i].append(verso[:spc_idx])
                new[i].append(verso[spc_idx+1:])
            else:
                new[i].append(verso)
    
    letra = new

    prs = get_title(prs, layout, artist, music, get_font_path(font_name), False)

    for estrofe in letra:
        for i, verso in enumerate(estrofe):
            add_slide = True if i % stacks == 0 else False
            times = 0 if i == (len(estrofe) - 1) else stacks - 1 - i % stacks
            prs = create_pwp(prs, layout, verso.upper(), font_name, size, \
                times=times, color=RGBColor(255, 255, 255), shadow=shadow, add_slide=add_slide)
    

    path = os.getcwd().replace("\\", "/") + "/letras"
    name = f"{path}/{music} - {artist}.pptx"

    return prs


def get_title(prs, layout, artist, music, font_name, shadow=False):
    prs = create_pwp(prs, layout, music, font_name, 32, times=0, \
        color=RGBColor(255, 255, 255), shadow=shadow, format="bold")
    prs = create_pwp(prs, layout, "Autor: " + artist, font_name, 32, \
        times=0, color=RGBColor(255, 255, 255), shadow=shadow, format="italic")
    

    # cons = 7.5 - 1 - 32 * .016 * (1 + 1)

    # slide = prs.slides[-1]

    # left = 0
    # top = Inches(cons)
    # width = prs.slide_width
    # height = Inches(.4)

    # music = "\"" + music + "\""

    # if shadow:
    #     txBox = slide.shapes.add_textbox(left - Inches(.015), top + Inches(.015), width, height)
    #     tf = txBox.text_frame
    #     p = tf.paragraphs[0]
    #     p.alignment = PP_ALIGN.CENTER

    #     run = p.add_run()
    #     # run.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    #     run.text = music

    #     font = run.font
    #     font.name = font_name
    #     font.size = Pt(32)
    #     font.color.rgb = RGBColor(0, 0, 0)

    # txBox = slide.shapes.add_textbox(left, top, width, height)
    # tf = txBox.text_frame
    # p = tf.paragraphs[0]

    # run = p.add_run()
    # # run.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    # run.text = music

    # font = run.font
    # font.name = font_name
    # font.size = Pt(32)
    # font.color.rgb = RGBColor(255, 255, 255)

    return prs


def init_slides():
    INDEX = 6
    prs = Presentation()

    # 13.333
    prs.slide_width = Inches(13.333)
    # 7.5
    prs.slide_height = Inches(7.5)

    # O index representa o tipo de layout que vamos usar
    layout = prs.slide_layouts[INDEX]

    return prs, layout


def create_pwp(prs, layout, text, font_name, size, times=0, \
    color=RGBColor(255, 255, 255), shadow=False, add_slide=True, format=None):
    """
    Cria Slides no Power Point de acordo com os inputs
    """
    # times = count_chars(text, "\n")
    border = 0.4
    cons = 7.4 - border - size * .016 * (times + 1)
    
    if add_slide:
        slide = prs.slides.add_slide(layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(32, 56, 100)
    else:
        slide = prs.slides[-1]

    # left = Inches(1)
    # # Adapta a distância do topo ao tamanho do texto
    # top = Inches(cons)
    # width = Inches(9)
    # height = Inches(.2 + .13 * times + .2 * (times+1))

    left = 0
    top = Inches(cons)
    width = prs.slide_width
    height = Inches(.2 + .13 * times + .2 * (times+1))

    if shadow:
        txBox = slide.shapes.add_textbox(left - Inches(.015), top + Inches(.015), width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER

        run = p.add_run()
        run.text = text

        font = run.font
        font.name = font_name
        font.size = Pt(size)
        font.color.rgb = RGBColor(0, 0, 0)
        if format == "bold":
            font.bold = True
        elif format == "italic":
            font.italic = True

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    run = p.add_run()
    run.text = text

    font = run.font
    font.name = font_name
    font.size = Pt(size)
    font.color.rgb = RGBColor(255, 255, 255)
    if format == "bold":
        font.bold = True
    elif format == "italic":
        font.italic = True

    # imgLeft = Inches(.8)
    # imgTop = Inches(.65)
    # imgWidth = Inches(1)
    # imgHeight = Inches(.838)
    # # imgHeight = Inches(imgWidth * .838)

    # img_path = "images/Logo da Juventude.png"
    # slide.shapes.add_picture(img_path, imgLeft, imgTop, imgWidth, imgHeight)

    # imgLeft = Inches(13.333 - 2.7)
    # imgTop = Inches(.65)
    # imgWidth = Inches(1 * 1.7)
    # imgHeight = Inches(.512 * 1.7)
    # # imgHeight = Inches(imgWidth * .838)

    # img_path = "images/Logo JMTV.png"
    # slide.shapes.add_picture(img_path, imgLeft, imgTop, imgWidth, imgHeight)

    return prs


def save(name, prs):
    path = name[:-name[::-1].index("/")]
    print("Path:", path)
    files = os.listdir(path)

    if name in files:
        os.remove(name)
    prs.save(name)
    # Open the file on the operating system
    # os.startfile(name)


# def count_chars(string, char):
#     counter = 0
#     c = len(char)
#     for i in range(0, len(string) - c, c):
#         if string[i:i+c] == char:
#             counter += 1
    
#     return counter


if __name__ == "__main__":
    main()
