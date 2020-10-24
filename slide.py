import os
import simplify
import tarfile
import shutil
import re
import io

from get_lyrics import requestsScrape, divide_by_text
import simplify

import matplotlib.colors
from PIL import Image, ImageDraw, ImageFont

from pptx import Presentation
from pptx.util import Pt, Inches, Cm
from pptx.util import Length
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN, MSO_ANCHOR


# def find(s, ch):
#     return (i for i, ltr in enumerate(s) if ltr == ch)


def split_sentence(sentence, times):
    """Splits the sentence in ``times`` parts"""
    idxs = [int((i+1) * (len(sentence) / times)) for i in range(times-1)]
    spc_idxs = [0]
    for i, char in enumerate(sentence):
        if char == " ": 
            spc_idxs.append(i)
    spc_idxs.append(None)

    splited = [] 

    if len(spc_idxs) - 2 < times:
        raise Exception("\"times\" must be less or equal the number of blank spaces!")
    elif len(spc_idxs) - 2 == times:
        for i, idx in enumerate(spc_idxs):
            if idx == None:
                break
            next = spc_idxs[i+1]
            splited.append(sentence[idx:next])
    else:
        difference = 100
        closest_idx = 0
        idxs_mapping = [0]

        spc_idxs[-1] = len(spc_idxs)

        for i, idx in enumerate(idxs):
            for j, spc_idx in enumerate(spc_idxs):
                if abs(spc_idx - idx) < difference:
                    difference = abs(spc_idx - idx)
                    closest_idx = spc_idx
            idxs_mapping.append(closest_idx)
            spc_idxs.remove(closest_idx)
            difference = 100

        idxs_mapping = sorted(idxs_mapping)
        idxs_mapping.append(None)

        for i, idx in enumerate(idxs_mapping):
            if idx == None:
                break
            next = idxs_mapping[i+1]
            splited.append(sentence[idx:next])
            
    return [i.strip() for i in splited]


class Px(Length):
    """
    Convenience value class for specifying a length in pixels
    """

    def __new__(cls, px):
        return Inches(px / 100)


class Colors(object):
    """
    Its attributes are color names and their values are rgb tuples
    """
    def __init__(self):
        for name, hex in matplotlib.colors.cnames.items():
            color = tuple([int(i * 255) for i in matplotlib.colors.to_rgb(hex)])
            setattr(self, name, color)


class Positions(object):
    """
    Maps all possible positions
    """
    def __init__(self, width, height, w, h, shadow=False):
        self.width = width
        self.height = height
        self.w = w
        self.h = h

        self.shadow = shadow

        self.border = 80
        self.extra = 5

        self.update()
    
    def update(self):
        top = self.border - self.extra * self.shadow
        bottom = self.height - self.border - self.h - self.extra * self.shadow
        h_center = int((self.height - self.h) / 2 - self.extra * self.shadow)
        w_center = int((self.width - self.w) / 2) - self.extra * self.shadow
        right = self.width - self.w - self.border * (self.width / self.height) - self.extra * self.shadow
        left = self.border * (self.width / self.height) - self.extra * self.shadow

        self.middle = (w_center, h_center)
        self.center_left = (left, h_center)
        self.center_right = (right, h_center)

        self.bottom_left = (left, bottom)
        self.bottom_right = (right, bottom)
        self.bottom_center = (w_center, bottom)

        self.top_left = (left, top)
        self.top_right = (right, top)
        self.top_center = (w_center, top)

    
    @staticmethod
    def get_font_size(text, font_type):
        if "\n" in text:
            greater = 0
            h = 0
            for sentence in text.split("\n"):
                size = font_type.getsize(sentence)
                w = size[0]
                greater = w if w > greater else greater

                # h += font_type.getsize("A")[1]
                h += font_type.getsize(sentence)[1]

            w = greater
        else:
            w, h = font_type.getsize(text)

        return w, h
    

    @staticmethod
    def check_alignment(font_name):
        translate_alignment = {
            'middle': 'center',
            'center_left': 'left',
            'center_right': 'right',
            'bottom_left': 'left',
            'bottom_right': 'right',
            'bottom_center': 'center',
            'top_left': 'left',
            'top_right': 'right',
            'top_center': 'center',
        }

        if font_name in translate_alignment:
            return translate_alignment[font_name]
        else:
            raise Exception("Invalid position!")


class FontType(object):
    def __init__(self, family_font, size, text_format=None):
        self._family_font = family_font
        self._size = size
        self._text_format = text_format

    def get_path(self):
        return simplify.get_font_path(self._family_font, self._text_format)
    
    def get_size(self):
        return self._size
    
    def get_family(self):
        return self._family_font
    
    def get_text_format(self):
        return self._text_format

    def set_size(self, size):
        self._size = size
    
    def set_family(self, family):
        self._family_font = family
    
    def set_text_format(self, text_format):
        self._text_format = text_format


class Slide(object):
    def __init__(self, background_path):
        self.width = 1920
        self.height = 1080

        self.background_path = background_path
        self.image_bytes = None

        self.font_type = FontType("Ubuntu", 40)
        self.shadow = False

        self.color = getattr(Colors(), "white")

        self.stacks = 2 # Number of sentences per slide

        self.directory = "../slides/Slides"

        self.border = 20

        self.extra = 5 # Shadow distance

        self.position = "middle"

        self.prs = None
    

    def adapt_size(self, lyrics, limit):
        """
        Returns a new list with the verses that are larger then the limit splited
        """
        new = []
        for i, estrofe in enumerate(lyrics):
            new.append([])
            for j, verso in enumerate(estrofe):
                font_type = ImageFont.truetype(self.font_type.get_path(), self.font_type.get_size())
                text_width = Positions.get_font_size(verso, font_type)[0]

                if text_width > limit:
                    split_times = text_width // limit

                    for subverso in split_sentence(verso, split_times):
                        new[i].append(subverso)
                else:
                    new[i].append(verso)
        
        return new
    
    def create_imageshow(self, lyrics):
        """
        Creates an image show
        """
        if not self.image_bytes:
            image = simplify.assign_image(self.background_path)
        else:
            image = simplify.assign_image(io.BytesIO(self.image_bytes))

        os.mkdir(self.directory)

        font_type = ImageFont.truetype(self.font_type.get_path(), self.font_type.get_size())

        normal_positions = None
        shadow_positions = None

        LIMIT = self.width - self.border * 2

        count = 1
        for estrofe in self.adapt_size(lyrics, LIMIT):
            for i in range(0, len(estrofe), self.stacks):
                if i < len(estrofe) - (self.stacks - 1):
                    text = "\n".join([estrofe[j].split()[0].capitalize() 
                                    + " " + " ".join(estrofe[j].split()[1:]) for j in range(i, i+self.stacks)])
                else:
                    last = len(estrofe)
                    text = "\n".join([estrofe[j].split()[0].capitalize() 
                                    + " " +  " ".join(estrofe[j].split()[1:]) for j in range(i, last)])
                
                w, h = Positions.get_font_size(text, font_type)

                if not normal_positions or not shadow_positions:
                    normal_positions = Positions(self.width, self.height, w, h)
                    shadow_positions = Positions(self.width, self.height, w, h, shadow=True)
                else:
                    normal_positions.w, normal_positions.h = w, h
                    shadow_positions.w, shadow_positions.h = w, h

                normal_positions.extra = self.extra
                normal_positions.border = self.border

                shadow_positions.extra = self.extra
                shadow_positions.border = self.border

                normal_positions.update()
                shadow_positions.update()

                normal_position = getattr(normal_positions, self.position)
                shadow_position = getattr(shadow_positions, self.position)

                image = simplify.assign_image(self.background_path)
                if self.shadow:
                    image = image.copy()
                    draw = ImageDraw.Draw(image)
                    draw.text(xy=shadow_position, text=text, fill=getattr(Colors(), "black"), \
                        font=font_type, align=Positions.check_alignment(self.position))

                image = image.copy()
                draw = ImageDraw.Draw(image)
                draw.text(xy=normal_position, text=text, fill=self.color, \
                    font=font_type, align=Positions.check_alignment(self.position))
                
                image.save(f"{self.directory}/{count}.png")

                count += 1
    
    def create_pwp(self, text):
        """
        Creates a pptx presentation based on the music lyrics
        """

        font_type = ImageFont.truetype(self.font_type.get_path(), self.font_type.get_size())
        w, h = Positions.get_font_size(text, font_type)

        normal_positions = Positions(self.width, self.height, w, h)
        shadow_positions = Positions(self.width, self.height, w, h, shadow=True)

        normal_positions.border = self.border
        
        shadow_positions.extra = self.extra
        shadow_positions.border = self.border

        normal_positions.update()
        shadow_positions.update()

        normal_position = getattr(normal_positions, self.position)
        shadow_position = getattr(shadow_positions, self.position)

        layout = self.prs.slide_layouts[6]
        
        slide = self.prs.slides.add_slide(layout)

        if not self.image_bytes:
            slide.shapes.add_picture(self.background_path, 0, 0, height = Px(self.height))
        else:
            slide.shapes.add_picture(io.BytesIO(self.image_bytes), 0, 0, height = Px(self.height))

        alignment = Positions.check_alignment(self.position)

        if "_" in self.position:
            h_alignment = self.position.split("_")[0].replace("center", "middle").upper()
        else:
            h_alignment = self.position.upper()

        if self.shadow:
            txBox = slide.shapes.add_textbox(Px(-self.extra), Px(-self.extra), 
                                             Px(self.width), Px(self.height))

            tf = txBox.text_frame
            tf.vertical_anchor = getattr(MSO_ANCHOR, h_alignment)

            if len(self.position.split("_")) > 1 and self.position.split("_")[1] == "left":
                tf.margin_left = Px(self.border)
            elif len(self.position.split("_")) > 1 and self.position.split("_")[1] == "right":
                tf.margin_right = Px(self.border)
            
            if len(self.position.split("_")) > 1 and self.position.split("_")[0] == "top":
                tf.margin_top = Px(self.border)
            elif len(self.position.split("_")) > 1 and self.position.split("_")[1] == "bottom":
                tf.margin_bottom = Px(self.border)

            tf.word_wrap = True

            p = tf.paragraphs[0]
            p.text = text
            p.alignment = getattr(PP_ALIGN, alignment.upper())

            font = p.font

            font.name = self.font_type.get_family()
            font.size = Px(self.font_type.get_size())
            font.color.rgb = RGBColor(*getattr(Colors(), "black"))

            font.bold = True if self.font_type.get_text_format() == "bold" else False
            font.italic = True if self.font_type.get_text_format() == "italic" else False

        txBox = slide.shapes.add_textbox(0, 0, Px(self.width), Px(self.height))

        tf = txBox.text_frame
        tf.vertical_anchor = getattr(MSO_ANCHOR, h_alignment)

        if len(self.position.split("_")) > 1 and self.position.split("_")[1] == "left":
            tf.margin_left = Px(self.border)
        elif len(self.position.split("_")) > 1 and self.position.split("_")[1] == "right":
            tf.margin_right = Px(self.border)
        
        if len(self.position.split("_")) > 1 and self.position.split("_")[0] == "top":
            tf.margin_top = Px(self.border)
        elif len(self.position.split("_")) > 1 and self.position.split("_")[0] == "bottom":
            tf.margin_bottom = Px(self.border)

        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = text
        p.alignment = getattr(PP_ALIGN, alignment.upper())

        font = p.font
        
        font.name = self.font_type.get_family()
        font.size = Px(self.font_type.get_size())
        font.color.rgb = RGBColor(*self.color)

        font.bold = True if self.font_type.get_text_format() == "bold" else False
        font.italic = True if self.font_type.get_text_format() == "italic" else False
    

    def create_slideshow(self, lyrics, font_format=[]):
        self.prs.slide_width = Px(self.width)
        self.prs.slide_height = Px(self.height)

        LIMIT = self.width - self.border * 2
        for i, estrofe in enumerate(self.adapt_size(lyrics, LIMIT)):
            for i in range(0, len(estrofe), self.stacks):
                if i < len(estrofe) - (self.stacks - 1):
                    text = "\n".join([estrofe[j].split()[0].capitalize() 
                                    + " " + " ".join(estrofe[j].split()[1:]) for j in range(i, i+self.stacks)])
                else:
                    last = len(estrofe)
                    text = "\n".join([estrofe[j].split()[0].capitalize() 
                                    + " " +  " ".join(estrofe[j].split()[1:]) for j in range(i, last)])

                # font_type = ImageFont.truetype(self.font_type.get_path(), self.font_type.get_size())

                self.create_pwp(text)
